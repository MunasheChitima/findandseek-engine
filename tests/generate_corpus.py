#!/usr/bin/env python3
"""Generate the test corpus in ``tests/corpus`` — the fixtures the ingest and
search tests run against.

Not a test (no ``test_`` prefix, so pytest ignores it) and not something users
run. ``tests/conftest.py`` calls it automatically when the corpus is missing, so
in normal use you never invoke it yourself:

    python tests/generate_corpus.py     # rewrites tests/corpus in place

**This script is the corpus.** The fixtures it writes are gitignored, not
committed, which makes it the single source of truth for what the tests run
against — and makes drift impossible by construction. That matters: the
artifacts *had* drifted. This script was changed to bill a fictional customer,
but the PDFs were never rebuilt, so a real person's name survived inside
`invoice_plumber_march.pdf` for as long as it was committed — invisible to every
text search, because the file is binary.

So: change a string here and it takes effect on the next test run. There is
nothing to remember to rebuild, and nothing stale to commit.

Every name, amount, and address below is invented. The README and the benchmark
docs both state the corpus is synthetic; ``tests/test_corpus_synthetic.py``
enforces it against this script's output.
"""

from __future__ import annotations

import csv
import json
import shutil
import zipfile
from io import BytesIO
from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageDraw, ImageFont

CORPUS = Path(__file__).resolve().parent / "corpus"


def write_invoice_pdf(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page()
    text = (
        "INVOICE\n\n"
        "Joe's Plumbing Services\n"
        "Invoice #: INV-2024-0312\n\n"
        "Bill To: Alex Rivera\n"
        "Date: March 15, 2024\n\n"
        "Description: Emergency pipe repair\n"
        "Amount Due: $485.00\n\n"
        "Payment due within 14 days."
    )
    page.insert_text((72, 72), text, fontsize=11)
    doc.save(str(path))
    doc.close()


def write_contract_pdf(path: Path) -> None:
    doc = fitz.open()
    paragraph = (
        "WHEREAS the parties herein agree to the following terms and conditions. "
        "The Contractor shall indemnify the Client against all claims arising from breach. "
        "Notice period shall not be less than four weeks for employees with over two years service. "
        "The Employee acknowledges receipt of this policy and agrees to abide by all provisions herein. "
        "Redundancy entitlements shall be calculated in accordance with the Fair Work Act and applicable awards. "
    )
    for i in range(30):
        page = doc.new_page()
        rect = fitz.Rect(72, 72, 520, 750)
        page.insert_textbox(rect, f"Section {i + 1}.\n\n{paragraph * 6}", fontsize=9)
    doc.save(str(path))
    doc.close()


def write_scanned_letter_pdf(path: Path) -> None:
    """Image-only PDF — text exists only as pixels."""
    img = Image.new("RGB", (600, 400), "white")
    draw = ImageDraw.Draw(img)
    draw.text((40, 40), "Dear Council,\n\nWe write regarding the fence dispute\non Oak Street.\n\nSigned, Clancy", fill="black")
    buf = BytesIO()
    img.save(buf, format="JPEG")
    doc = fitz.open()
    page = doc.new_page(width=600, height=400)
    page.insert_image(page.rect, stream=buf.getvalue())
    doc.save(str(path))
    doc.close()


def write_board_deck(path: Path) -> None:
    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[1])
    title.shapes.title.text = "Board Deck Q3 2025"
    title.placeholders[1].text = "Revenue overview and KPIs"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    img = Image.new("RGB", (400, 200), "white")
    draw = ImageDraw.Draw(img)
    draw.text((20, 80), "Q3 Revenue: $1.2M (+15% YoY)", fill="black")
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(1), Inches(1), width=Inches(5))
    prs.save(str(path))


def write_cashflow_xlsx(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Cashflow"
    ws.append(["date", "category", "amount", "balance"])
    for m in range(1, 13):
        ws.append([f"2025-{m:02d}-01", "operations", 1000 * m, 5000 + m * 100])
    ws2 = wb.create_sheet("Summary")
    ws2.append(["metric", "value"])
    ws2.append(["total", 78000])
    wb.save(str(path))


def write_contacts_csv(path: Path) -> None:
    with path.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["name", "email", "phone"])
        w.writerow(["Alice Smith", "alice@example.com", "+1 555-0100"])
        w.writerow(["Bob Jones", "bob@example.com", "+1 555-0101"])


def write_letter_docx(path: Path) -> None:
    doc = Document()
    doc.add_paragraph("Letter to the Council")
    doc.add_paragraph(
        "Dear Council Members,\n\n"
        "I am writing on behalf of residents regarding the proposed park renovation. "
        "Clancy Morrison spoke at the March 12, 2025 meeting and raised concerns about noise. "
        "We respectfully request a community consultation before work begins."
    )
    doc.save(str(path))


def write_notes(path: Path) -> None:
    path.write_text("# Project Notes\n\nFindandSeek test fixture.\nRemember to index Downloads.\n", encoding="utf-8")


def write_email(path: Path) -> None:
    raw = """From: sender@example.com
To: recipient@example.com
Subject: Re: Budget approval
Date: Mon, 3 Jun 2024 10:00:00 +0000
Content-Type: text/plain; charset=utf-8

Please review the attached budget spreadsheet.
Attachment: budget_q2.xlsx
"""
    path.write_text(raw, encoding="utf-8")


def write_data_json(path: Path) -> None:
    path.write_text(json.dumps({"project": "findandseek", "version": 1, "tags": ["local", "search"]}, indent=2), encoding="utf-8")


def write_photo_sign(path: Path) -> None:
    img = Image.new("RGB", (300, 150), "lightgray")
    draw = ImageDraw.Draw(img)
    draw.text((20, 60), "STOP — Private Property", fill="red")
    img.save(str(path))


def write_corrupt_pdf(path: Path) -> None:
    path.write_bytes(b"%PDF-1.4\n%corrupt truncated file")


def write_empty_docx(path: Path) -> None:
    Document().save(str(path))


def write_huge_text(path: Path) -> None:
    paragraph = "Lorem ipsum dolor sit amet. " * 50 + "\n"
    path.write_text(paragraph * 800, encoding="utf-8")


def main() -> None:
    if CORPUS.exists():
        shutil.rmtree(CORPUS)
    CORPUS.mkdir(parents=True)

    write_invoice_pdf(CORPUS / "invoice_plumber_march.pdf")
    write_contract_pdf(CORPUS / "contract_long.pdf")
    write_scanned_letter_pdf(CORPUS / "scanned_letter.pdf")
    write_board_deck(CORPUS / "board_deck.pptx")
    write_cashflow_xlsx(CORPUS / "cashflow.xlsx")
    write_contacts_csv(CORPUS / "contacts.csv")
    write_letter_docx(CORPUS / "letter_to_council.docx")
    write_notes(CORPUS / "notes.md")
    write_notes(CORPUS / "notes.txt")
    write_email(CORPUS / "email_thread.eml")
    write_data_json(CORPUS / "data.json")
    write_photo_sign(CORPUS / "photo_with_sign.jpg")
    write_corrupt_pdf(CORPUS / "corrupt.pdf")
    shutil.copy(CORPUS / "invoice_plumber_march.pdf", CORPUS / "duplicate_of_1.pdf")
    write_empty_docx(CORPUS / "empty.docx")
    write_huge_text(CORPUS / "huge_text.txt")
    print(f"Generated corpus at {CORPUS}")


if __name__ == "__main__":
    main()
