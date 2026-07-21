"""The test corpus must contain no real personal data.

The README and the benchmark docs both state that every fixture is invented.
This test makes that checkable in CI instead of taking it on trust.

Why it exists: the committed fixtures drifted from `generate_corpus.py`. The
generator was changed to bill a fictional "Alex Rivera", but the PDFs were never
rebuilt, so the maintainer's real name survived inside
`invoice_plumber_march.pdf` — invisible to every text search, because the file
is binary. A byte-hash check can't catch this (PDF/docx/xlsx embed build
timestamps, so generation isn't reproducible), but a *content* check can.

The load-bearing assertion is `test_maintainer_is_not_a_data_subject`: whoever
maintains the project should never appear inside a fixture. That generalises
past the one name that actually leaked.
"""

from __future__ import annotations

import json
import re
import tomllib
import warnings
from pathlib import Path

import pytest

CORPUS = Path(__file__).resolve().parent / "corpus"
PYPROJECT = Path(__file__).resolve().parents[1] / "pyproject.toml"


def extract_text(path: Path) -> str:
    """Best-effort text for any fixture format, including the binary ones."""
    suffix = path.suffix.lower()
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        try:
            if suffix == ".pdf":
                import fitz

                with fitz.open(path) as doc:
                    return "".join(page.get_text() for page in doc)
            if suffix == ".docx":
                import docx

                return "\n".join(p.text for p in docx.Document(path).paragraphs)
            if suffix == ".xlsx":
                import openpyxl

                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
                return "\n".join(
                    str(c)
                    for ws in wb
                    for row in ws.iter_rows(values_only=True)
                    for c in row
                    if c is not None
                )
            if suffix == ".pptx":
                from pptx import Presentation

                return "\n".join(
                    shape.text
                    for slide in Presentation(path).slides
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                )
            return path.read_text(errors="ignore")
        except Exception:  # noqa: BLE001 - an unreadable fixture is corrupt.pdf, by design
            return ""


def corpus_files() -> list[Path]:
    return sorted(p for p in CORPUS.iterdir() if p.is_file())


@pytest.fixture(scope="module")
def corpus_text() -> dict[str, str]:
    return {p.name: extract_text(p) for p in corpus_files()}


def test_corpus_exists():
    assert CORPUS.is_dir(), "run `python tests/generate_corpus.py`"
    assert len(corpus_files()) >= 10


def test_maintainer_is_not_a_data_subject(corpus_text):
    """The project's own author must not appear inside a fixture.

    Being named in LICENSE.md, pyproject, and commit metadata is expected. Being
    the "Bill To:" on a sample invoice is not — that means a real document was
    used as a fixture, or a real name was typed into a synthetic one.
    """
    authors = tomllib.loads(PYPROJECT.read_text())["project"].get("authors", [])
    names = [a["name"] for a in authors if a.get("name")]
    assert names, "pyproject declares no author — nothing to check against"

    tokens = {t.lower() for n in names for t in n.split() if len(t) > 2}
    offenders = {
        name: sorted(t for t in tokens if re.search(rf"\b{re.escape(t)}\b", text, re.I))
        for name, text in corpus_text.items()
    }
    offenders = {k: v for k, v in offenders.items() if v}
    assert not offenders, (
        f"maintainer name found in fixtures: {offenders}. "
        "Fix the string in tests/generate_corpus.py and re-run it — do not edit "
        "the binary artifacts by hand."
    )


def test_no_contact_details(corpus_text):
    """No real-looking email addresses or phone numbers.

    example.com / example.org are reserved for documentation (RFC 2606), so
    fixtures may use them freely.
    """
    email = re.compile(r"\b[\w.+-]+@(?!example\.(?:com|org|net))[\w-]+\.[a-z]{2,}\b", re.I)
    phone = re.compile(r"\b(?:\+?61\s?4|04)\d{2}\s?\d{3}\s?\d{3}\b")

    found = {
        name: sorted(set(email.findall(text)) | set(phone.findall(text)))
        for name, text in corpus_text.items()
    }
    found = {k: v for k, v in found.items() if v}
    assert not found, f"contact details in fixtures: {found}"


def test_no_identifiers_from_a_local_denylist(corpus_text):
    """Optional per-machine denylist, for terms that must never reach a fixture.

    Deliberately **not** enumerated in this file. An explicit list of "strings
    that once leaked from this repository" is itself a disclosure: published, it
    tells a reader exactly what to look for, which for health information or a
    named third party is worse than the original mistake. So the terms live in
    an untracked file that never leaves the machine that needs it:

        tests/.denylist        # one term per line, '#' comments allowed

    Absent — the normal case, and the case for anyone cloning — this skips. The
    general guards above (maintainer name, contact details) are what actually
    run in CI, and they are the ones that catch the realistic regression.
    """
    denylist_path = Path(__file__).parent / ".denylist"
    if not denylist_path.exists():
        pytest.skip("no tests/.denylist on this machine")

    terms = [
        line.strip()
        for line in denylist_path.read_text().splitlines()
        if line.strip() and not line.lstrip().startswith("#")
    ]
    if not terms:
        pytest.skip("tests/.denylist is empty")

    hits = {
        name: [t for t in terms if re.search(rf"\b{re.escape(t)}\b", text, re.I)]
        for name, text in corpus_text.items()
    }
    hits = {k: v for k, v in hits.items() if v}
    assert not hits, f"denylisted terms present in fixtures: {sorted(hits)}"


def test_gold_eval_set_targets_only_this_corpus():
    """Every gold query must expect a file that this corpus actually contains.

    Stronger than a denylist and it names nothing: the eval set is queries over
    the synthetic corpus, so any `expect` substring that matches no generated
    fixture means the set is pointed at documents from somewhere else — which is
    how a personal index ends up referenced in a public repository.
    """
    gold = Path(__file__).resolve().parents[1] / "find_and_seek" / "eval" / "gold.jsonl"
    if not gold.exists():
        pytest.skip("no gold set")

    stems = {p.stem.lower() for p in corpus_files()} | {p.name.lower() for p in corpus_files()}
    dangling: dict[str, list[str]] = {}
    for line in gold.read_text().splitlines():
        if not line.strip():
            continue
        entry = json.loads(line)
        missing = [
            e for e in entry.get("expect", [])
            if not any(e.lower() in stem for stem in stems)
        ]
        if missing:
            dangling[entry.get("id", "?")] = missing

    assert not dangling, (
        f"gold queries expect files this corpus does not contain: {dangling}. "
        "The eval set must target tests/corpus, not a personal index."
    )
