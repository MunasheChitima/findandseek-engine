"""PPTX per-slide text + speaker notes + images."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from find_and_seek.ingest.extract.base import ExtractResult, ImageRegion, TextBlock


def _slide_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(f"[Notes] {notes}")
    return "\n".join(parts)


def extract_pptx(path: Path) -> ExtractResult:
    prs = Presentation(str(path))
    blocks: list[TextBlock] = []
    images: list[ImageRegion] = []

    for idx, slide in enumerate(prs.slides, start=1):
        text = _slide_text(slide)
        loc = f"slide {idx}"
        if text:
            blocks.append(TextBlock(text=text, source_type="text", location_ref=loc))

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    images.append(
                        ImageRegion(
                            data=blob,
                            mime="image/png",
                            location_ref=f"slide {idx} image",
                        )
                    )
                except (AttributeError, ValueError):
                    continue

    return ExtractResult(
        blocks=blocks,
        images=images,
        page_count=len(prs.slides),
        file_type="presentation",
    )
